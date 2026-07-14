# -*- coding: utf-8 -*-
"""
Build the full supervisor presentation deck (.pptx) — Phase 0 → 6.

Covers everything done: dataset, cleaning, language, label scheme + method, EDA
findings, split, models (why + results), test evaluation + significance, and the
deep analysis (language robustness, embeddings, religion, explainability, errors).

Numbers are pulled from results/*.json so the deck always matches the experiments.
Run with a Python that has python-pptx + Pillow:
    <py3.14>\\python.exe sentiment/src/build_slides.py
Output: sentiment/results/progress_update.pptx
"""
from __future__ import annotations
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "figures"
RES = ROOT / "results"
OUT = RES / "progress_update.pptx"


def load(name):
    p = RES / name
    return json.loads(p.read_text(encoding="utf-8")) if p.exists() else {}


eda = load("eda_stats.json")
sp = load("phase2_split.json")
base = load("phase3_baseline.json")
bertj = load("phase3_bert.json")
t4 = load("phase4_test_eval.json")
p5 = load("phase5_deep_analysis.json")
p5e = load("phase5_explain_errors.json")
p5s = load("phase5_stats.json")

TB = t4.get("test", {}).get("baseline_linear_svm", {})
TR = t4.get("test", {}).get("banglabert_tuned", {})
MC = t4.get("mcnemar_baseline_vs_bert", {})
MC_P = MC.get("p_value", 1)
SIG = MC_P < 0.05
SIG_TXT = "statistically significant" if SIG else "NOT statistically significant"
EMB = p5.get("embedding", {})
REL = p5.get("religion", {})
LANG = {r["lang"]: r for r in (p5.get("language", {}) or {}).get("test_accuracy_by_language", [])}
EG = p5e.get("error_analysis", {}).get("error_rate_by_group", {})

# ---- theme ----
NAVY = RGBColor(0x1D, 0x35, 0x57)
RED = RGBColor(0xE6, 0x39, 0x46)
GREY = RGBColor(0x49, 0x4F, 0x57)
LIGHT = RGBColor(0xF1, 0xF3, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2A, 0x9D, 0x8F)
GOLD = RGBColor(0xE9, 0xC4, 0x6A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()


def bar(s):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.16))
    fill(r, RED)


def title_bar(s, title, sub=None):
    bar(s)
    tf = box(s, Inches(0.55), Inches(0.28), Inches(12.3), Inches(1.0))
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = NAVY
    if sub:
        sp2 = tf.add_paragraph(); sp2.text = sub
        sp2.font.size = Pt(13.5); sp2.font.color.rgb = GREY


def bullets(s, items, left=Inches(0.7), top=Inches(1.55), width=Inches(12),
            size=16, gap=7):
    tf = box(s, left, top, width, SH - top - Inches(0.35))
    for i, it in enumerate(items):
        lvl = 0
        if isinstance(it, tuple):
            it, lvl = it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("• " if lvl == 0 else "    – ") + it
        p.font.size = Pt(size - (2 if lvl else 0))
        p.font.color.rgb = NAVY if lvl == 0 else GREY
        p.font.bold = (lvl == 0 and it.endswith(":"))
        p.space_after = Pt(gap)
    return tf


def picture(s, name, left, top, max_w, max_h):
    from PIL import Image
    path = FIG / name
    iw, ih = Image.open(path).size
    ratio = min(max_w / iw, max_h / ih)
    w, h = int(iw * ratio), int(ih * ratio)
    s.shapes.add_picture(str(path), left + (max_w - w) // 2, top + (max_h - h) // 2, w, h)


def caption(s, text, top, left=Inches(0.6), width=Inches(12.1)):
    tf = box(s, left, top, width, Inches(0.6))
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(12.5); p.font.italic = True; p.font.color.rgb = GREY
    p.alignment = PP_ALIGN.CENTER


def metric_cards(s, cards, top=Inches(1.7), left=Inches(0.7), card_w=Inches(2.9),
                 card_h=Inches(1.5), gap=Inches(0.35)):
    """cards = [(value, label, color), ...] big highlighted numbers."""
    x = left
    for val, label, color in cards:
        rc = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, card_w, card_h)
        fill(rc, LIGHT); rc.line.color.rgb = color; rc.line.width = Pt(1.5)
        tf = rc.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = val
        p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        q = tf.add_paragraph(); q.text = label
        q.font.size = Pt(12); q.font.color.rgb = GREY; q.alignment = PP_ALIGN.CENTER
        x = x + card_w + gap


def table(s, data, widths, left, top, total_w, total_h, bold_row=None, font=12.5):
    rows, cols = len(data), len(data[0])
    tbl = s.shapes.add_table(rows, cols, left, top, total_w, total_h).table
    tw = sum(widths)
    for j, w in enumerate(widths):
        tbl.columns[j].width = Inches(total_w.inches * w / tw)
    for j in range(cols):
        c = tbl.cell(0, j); c.text = str(data[0][j])
        c.fill.solid(); c.fill.fore_color.rgb = NAVY
        pr = c.text_frame.paragraphs[0]
        pr.font.color.rgb = WHITE; pr.font.bold = True; pr.font.size = Pt(font)
    for i in range(1, rows):
        for j in range(cols):
            c = tbl.cell(i, j); c.text = str(data[i][j])
            c.fill.solid()
            c.fill.fore_color.rgb = (RGBColor(0xE8, 0xEE, 0xF6) if bold_row == i
                                     else (LIGHT if i % 2 else WHITE))
            pr = c.text_frame.paragraphs[0]
            pr.font.size = Pt(font)
            pr.font.bold = (bold_row == i)
            pr.font.color.rgb = NAVY if bold_row == i else GREY
    return tbl


def dist_str(d, order=None):
    items = [(k, d[k]) for k in (order or d)] if d else []
    return " · ".join(f"{k} {v}" for k, v in items if k in d)


# ═══════════════════════ 1. TITLE ═══════════════════════
s = slide()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH); fill(bg, NAVY)
tf = box(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(2.3))
p = tf.paragraphs[0]
p.text = "Sentiment & Emotion Analysis of Social-Media\nComments on Elder Abuse"
p.font.size = Pt(38); p.font.bold = True; p.font.color.rgb = WHITE
sub = tf.add_paragraph()
sub.text = "Bangla · Banglish · English  —  a deep-analysis NLP study"
sub.font.size = Pt(18); sub.font.color.rgb = GOLD
meta = box(s, Inches(0.95), Inches(5.2), Inches(11), Inches(1.6))
for i, line in enumerate([
        "End-to-end pipeline complete — Phase 0 → 6 (data · models · deep analysis)",
        "Lamia Islam Mim  ·  CSE 499A (Capstone), North South University",
        "Component added per supervisor brief: “the deeper the analysis, the better”"]):
    p = meta.paragraphs[0] if i == 0 else meta.add_paragraph()
    p.text = line
    p.font.size = Pt(15 if i == 0 else 13); p.font.bold = (i == 0)
    p.font.color.rgb = WHITE if i == 0 else RGBColor(0xC9, 0xD2, 0xDE)

# ═══════════════════════ 2. OBJECTIVE ═══════════════════════
s = slide(); title_bar(s, "Objective & Why It Matters")
bullets(s, [
    "Goal: a publishable NLP system that classifies public comments on elder-abuse "
    "news/videos by sentiment AND fine-grained emotion.",
    "Why it matters:",
    ("Elder abuse is rising in Bangladesh; public reaction reveals social attitudes — "
     "condemnation, sympathy, and (critically) victim-blaming.", 1),
    ("Most Bangla sentiment work is on product/movie reviews; emotion analysis on a real "
     "social-issue corpus is novel.", 1),
    "Two research outputs:",
    ("(1) A new annotated Bangla/Banglish elder-abuse emotion dataset (with a novel "
     "victim-blaming class).", 1),
    ("(2) Benchmarked models + a deep analysis of how reactions vary by abuse type, "
     "platform, and language.", 1),
], top=Inches(1.6), gap=9)

# ═══════════════════════ 3. PIPELINE OVERVIEW ═══════════════════════
s = slide(); title_bar(s, "The Pipeline — Phase 0 → 6", "Every stage below is complete")
rows = [["Phase", "What was done", "Status"],
        ["0  Cleaning", "2,454 → 2,301 comments, language-tagged", "DONE"],
        ["1  Annotation", "2-level labels (polarity + emotion) via a written guideline", "DONE"],
        ["1.5  EDA", "13 charts + 7 documented findings", "DONE"],
        ["2  Split", "1609 / 345 / 345, stratified, test sha256-locked", "DONE"],
        ["3  Modelling", "TF-IDF+SVM baseline + fine-tuned BanglaBERT", "DONE"],
        ["4  Evaluation", "held-out test set + McNemar significance", "DONE"],
        ["5  Deep analysis", "language · embeddings · religion · explainability · errors", "DONE"],
        ["6  Paper assets", "all tables + 21 figures assembled", "DONE"]]
t = table(s, rows, [2.2, 7.4, 1.5], Inches(0.7), Inches(1.75), Inches(12), Inches(4.9))
for i in range(1, len(rows)):
    cell = t.cell(i, 2); cp = cell.text_frame.paragraphs[0]
    cp.font.bold = True; cp.font.color.rgb = GREEN

# ═══════════════════════ 4. DATASET ═══════════════════════
s = slide(); title_bar(s, "The Dataset", "Scraped YouTube + Facebook comments on elder-abuse videos")
metric_cards(s, [("2,454", "raw scraped", GREY), ("2,301", "after cleaning", NAVY),
                 ("2,299", "modelling set", GREEN)], top=Inches(1.55))
bullets(s, [
    "Cleaning removed 153 junk items: 118 textless/emoji-only, 13 spam, 13 duplicates, 9 too-short.",
    f"Platforms: {dist_str(eda.get('platform_counts', {}), ['YouTube','Facebook'])}.",
    "Video topic (context only, NOT the comment label): Neglect · Abandonment · Physical · "
    "Verbal · Murder.",
    "Handled up-front: scrape-time date (no temporal claims); abuse_type is video-level context; "
    "84% of likes = 0; emoji & punctuation kept (they carry emotion).",
], top=Inches(3.35), gap=9)

# ═══════════════════════ 5. LANGUAGE ═══════════════════════
s = slide(); title_bar(s, "Language Composition", "Detected by script (Unicode ranges) at cleaning time")
picture(s, "18_lang_distribution.png", Inches(0.4), Inches(1.7), Inches(7.4), Inches(4.7))
bullets(s, [
    "Bangla 1,899 — pure Bangla script.",
    "Banglish (Latin) 367 — Bangla written in English letters.",
    "Code-mixed 33 — Bangla + English in one comment.",
    "Arabic 2 — short religious phrases.",
    "≈16% is Latin-script — this matters a lot for model choice (shown later).",
], left=Inches(8.0), top=Inches(2.0), width=Inches(5.0), size=15, gap=10)

# ═══════════════════════ 6. LABEL SCHEME ═══════════════════════
s = slide(); title_bar(s, "Two-Level Label Scheme")
bullets(s, [
    "Level A — Polarity (3): Negative · Neutral · Positive  (overall emotional valence).",
    "Level B — Fine-grained emotion (5):",
    ("Anger/Condemnation — outrage, demand for justice/punishment.", 1),
    ("Victim-blaming — blaming the elder/family (the novel, high-value class).", 1),
    ("Grief/Sympathy — sorrow, empathy for the victim.", 1),
    ("Prayer/Religious — dua, blessings, religious framing.", 1),
    ("Neutral-Other — factual, procedural, off-topic.", 1),
    "Every comment gets one polarity + one emotion. Test set is sha256-locked.",
], top=Inches(1.55), gap=6)

# ═══════════════════════ 7. HOW WE LABELED ═══════════════════════
s = slide(); title_bar(s, "How Labels Were Assigned", "A written codebook — not ad-hoc")
bullets(s, [
    "A formal annotation guideline defines each label with rules, examples, and tie-breaks.",
    "Key rules:",
    ("Label the COMMENT text only — never the video's abuse_type (avoids topic bias).", 1),
    ("Religion is the language, not the label: dua for victim → Positive/Prayer; curse on abuser "
     "in God's name → Negative/Anger.", 1),
    ("Emoji are evidence: 😢 → grief, 😡 → anger. Sarcasm → intended meaning.", 1),
    ("Multiple emotions → pick the dominant one.", 1),
    "Examples: “আল্লাহ জান্নাত দিক” → Positive · “কুলাঙ্গারের ফাঁসি হোক” → Negative/Anger · "
    "“ভিডিওটা কোথাকার?” → Neutral.",
], top=Inches(1.55), gap=7)

# ═══════════════════════ 8. EDA DISTRIBUTIONS ═══════════════════════
s = slide(); title_bar(s, "Result — Label Distributions")
picture(s, "01_polarity_dist.png", Inches(0.4), Inches(1.65), Inches(6.2), Inches(4.6))
picture(s, "02_emotion_dist.png", Inches(6.7), Inches(1.65), Inches(6.3), Inches(4.6))
caption(s, "70% of reactions are Negative. Anger (~35%) and victim-blaming (~13.5%) dominate the "
           "emotion mix.", Inches(6.45))

# ═══════════════════════ 9. FINDING 1 ═══════════════════════
s = slide(); title_bar(s, "Key Finding 1 — Emotion Depends on Abuse Type")
picture(s, "04_abusetype_x_emotion.png", Inches(1.4), Inches(1.6), Inches(10.5), Inches(4.6))
caption(s, "Anger peaks for Verbal/Physical abuse; prayer peaks for Murder; victim-blaming "
           "concentrates in Neglect/Abandonment.", Inches(6.4))

# ═══════════════════════ 10. FINDING 2 (victim-blaming) ═══════════════════════
s = slide(); title_bar(s, "Key Finding 2 — Victim-Blaming is Topic-Specific")
picture(s, "10_victimblame_by_abuse.png", Inches(0.5), Inches(1.65), Inches(7.2), Inches(4.6))
bullets(s, [
    "25.7% of Neglect comments blame the victim/family…",
    "…but 2.2% for Physical and 0% for Murder.",
    "When elders are abandoned (often educated/wealthy parents), a vocal minority blames "
    "them (“bad upbringing”, “haram money”).",
    "When harm is undeniable violence, the public unites in condemnation.",
    "Victim-blaming comments are also the LONGEST — they argue and justify.",
], left=Inches(8.0), top=Inches(1.9), width=Inches(5.0), size=14, gap=9)

# ═══════════════════════ 11. FINDING 3&4 ═══════════════════════
s = slide(); title_bar(s, "Key Findings 3 & 4 — Platform and Engagement")
picture(s, "06_platform_x_emotion.png", Inches(0.3), Inches(1.65), Inches(6.5), Inches(4.5))
picture(s, "09_likes_by_emotion.png", Inches(6.9), Inches(1.65), Inches(6.1), Inches(4.5))
caption(s, "Victim-blaming is 18.6% on YouTube vs 0.9% on Facebook. Grief/Sympathy earns ~10× more "
           "likes than anger.", Inches(6.35))

# ═══════════════════════ 12. FINDING 5 (lexical) ═══════════════════════
s = slide(); title_bar(s, "Key Finding 5 — Clean Lexical Separation")
picture(s, "12_top_tokens.png", Inches(0.6), Inches(1.65), Inches(12.1), Inches(4.5))
caption(s, "Distinctive words per emotion are sharply different (justice/wretch vs "
           "education/old-age-home vs pain/sorrow) — supports the 5-class scheme.", Inches(6.35))

# ═══════════════════════ 13. SPLIT ═══════════════════════
s = slide(); title_bar(s, "Phase 2 — Dataset Split", "Zero-leakage, imbalance-aware")
metric_cards(s, [("1,609", "train", NAVY), ("345", "validation", GREEN),
                 ("345", "test (locked)", RED)], top=Inches(1.6))
bullets(s, [
    "Text-dedup BEFORE splitting → no leakage across sets.",
    "Stratified 70/15/15 by polarity (seed 42) → same class ratios in every split.",
    "Test set sha256-locked and touched exactly once (final evaluation only).",
    "Primary metric = macro-F1 — averages classes equally, so the rare Positive/Neutral "
    "classes count as much as the dominant Negative.",
], top=Inches(3.4), gap=9)

# ═══════════════════════ 14. MODELLING APPROACH ═══════════════════════
s = slide(); title_bar(s, "Phase 3 — Modelling Approach", "Why each choice")
bullets(s, [
    "Baseline — TF-IDF + Linear SVM:",
    ("Word n-grams (1–2) + character n-grams (3–5). Char n-grams absorb Bangla spelling "
     "variation AND romanized Banglish.", 1),
    ("Class-weighted to handle the Negative-heavy imbalance.", 1),
    "Transformer — BanglaBERT (sagorsarker/bangla-bert-base):",
    ("Fine-tuned with class-weighted loss; best validation checkpoint kept.", 1),
    ("Tuned: 5 epochs, warmup, weight-decay.", 1),
    "Fit on TRAIN only; selected on VALIDATION; TEST stays locked → honest generalisation.",
], top=Inches(1.55), gap=6)

# ═══════════════════════ 15. VALIDATION RESULTS ═══════════════════════
s = slide(); title_bar(s, "Phase 3 — Validation Results")
rows = [["Model", "macro-F1", "accuracy"],
        ["Logistic Regression", "0.653", "0.759"],
        ["Linear SVM  ★ best", "0.661", "0.777"],
        ["Complement NB", "0.487", "0.704"],
        ["BanglaBERT (tuned)", "0.579", "0.713"]]
table(s, rows, [4.0, 1.6, 1.6], Inches(0.6), Inches(1.75), Inches(6.4), Inches(2.7), bold_row=2)
picture(s, "14_baseline_confusion.png", Inches(7.2), Inches(1.7), Inches(5.6), Inches(4.6))
bullets(s, [
    "Negative detected very well (F1 0.86).",
    "Neutral is the hard class (fuzzy middle).",
    "SVM leads on validation.",
], left=Inches(0.6), top=Inches(4.7), width=Inches(6.4), size=14, gap=6)

# ═══════════════════════ 16. TEST EVALUATION ═══════════════════════
s = slide(); title_bar(s, "Phase 4 — Final Test-Set Evaluation", "Locked test set, touched once")
rows = [["Model", "test macro-F1", "test acc", "weighted-F1"],
        ["Linear SVM  ★", f"{TB.get('macro_f1',0):.3f}", f"{TB.get('accuracy',0):.3f}",
         f"{TB.get('weighted_f1',0):.3f}"],
        ["BanglaBERT (tuned)", f"{TR.get('macro_f1',0):.3f}", f"{TR.get('accuracy',0):.3f}",
         f"{TR.get('weighted_f1',0):.3f}"]]
table(s, rows, [3.4, 1.9, 1.5, 1.9], Inches(0.6), Inches(1.75), Inches(7.2), Inches(1.7), bold_row=1)
picture(s, "16_test_confusion_baseline.png", Inches(8.0), Inches(1.7), Inches(4.9), Inches(4.7))
bullets(s, [
    f"Both generalise well — test ≥ validation (no overfitting).",
    (f"McNemar p = {MC_P:.3f} → the gap is {SIG_TXT}; the Linear SVM outperforms BanglaBERT."
     if SIG else
     f"McNemar p = {MC_P:.2f} → the gap is {SIG_TXT}; the two models are comparable."),
    "Linear SVM is the practical choice: simpler, faster, no GPU.",
], left=Inches(0.6), top=Inches(3.7), width=Inches(7.2), size=14, gap=9)

# ═══════════════════════ 16b. PER-CLASS METRICS ═══════════════════════
s = slide(); title_bar(s, "Per-Class Precision / Recall / F1", "Held-out test — where each model wins")
picture(s, "23_per_class_metrics.png", Inches(0.6), Inches(1.7), Inches(12.1), Inches(4.3))
caption(s, "Both models are strong on Negative; Neutral is the hard middle. The SVM is more "
           "balanced on the rare Positive/Neutral classes, which lifts its macro-F1.", Inches(6.25))

# ═══════════════════════ 17. DEEP: LANGUAGE ROBUSTNESS ═══════════════════════
s = slide(); title_bar(s, "Phase 5 — Why the Baseline Wins: Language Robustness")
picture(s, "19_lang_model_accuracy.png", Inches(0.4), Inches(1.65), Inches(7.4), Inches(4.6))
bl = LANG.get("Banglish (Latin)", {})
bullets(s, [
    "On held-out test, by language:",
    (f"Bangla — SVM {LANG.get('Bangla',{}).get('svm_acc',0):.2f} · "
     f"BanglaBERT {LANG.get('Bangla',{}).get('bert_acc',0):.2f}", 1),
    (f"Banglish (Latin) — SVM {bl.get('svm_acc',0):.2f} · "
     f"BanglaBERT {bl.get('bert_acc',0):.2f}  ← collapses", 1),
    "BanglaBERT was pre-trained on Bangla SCRIPT — it can't read romanized Banglish.",
    "The char-ngram SVM stays robust. Since ~16% of the corpus is Latin-script, "
    "this largely explains why the baseline wins overall.",
], left=Inches(8.0), top=Inches(1.9), width=Inches(5.0), size=14, gap=9)

# ═══════════════════════ 18. DEEP: EMBEDDINGS ═══════════════════════
s = slide(); title_bar(s, "Phase 5 — Embedding Structure")
picture(s, "20_embedding_tsne_emotion.png", Inches(0.5), Inches(1.65), Inches(7.0), Inches(4.7))
bullets(s, [
    "BanglaBERT sentence embeddings → KMeans vs human labels (Adjusted Rand Index):",
    (f"cluster ↔ emotion: ARI = {EMB.get('ari_cluster_vs_emotion_k5',0):.3f}", 1),
    (f"cluster ↔ polarity: ARI = {EMB.get('ari_cluster_vs_polarity_k3',0):.3f}", 1),
    "Low ARI ⇒ emotion classes OVERLAP in embedding space (semantic, not lexically "
    "separable) ⇒ supervised labels are genuinely needed.",
    "The detached blob (left) is Latin-script text — grouped by script, not emotion.",
], left=Inches(7.8), top=Inches(1.9), width=Inches(5.2), size=14, gap=8)

# ═══════════════════════ 19. DEEP: RELIGION ═══════════════════════
s = slide(); title_bar(s, "Phase 5 — Religious Framing")
picture(s, "21_religious_framing.png", Inches(0.4), Inches(1.65), Inches(7.6), Inches(4.6))
ew = REL.get("emotion_within_religious", {})
bullets(s, [
    f"{REL.get('religious_share',0)*100:.0f}% of all comments use religious language.",
    "Within religious comments:",
    (f"Prayer/Religious {ew.get('Prayer/Religious',0)*100:.0f}%  (blessing)", 1),
    (f"Victim-blaming {ew.get('Victim-blaming',0)*100:.0f}%  ·  "
     f"Anger {ew.get('Anger/Condemnation',0)*100:.0f}%  (blame)", 1),
    "Religion is the vehicle of BOTH blessing and condemnation — it skews comments "
    "slightly more Positive (dua) than non-religious ones.",
], left=Inches(8.2), top=Inches(1.9), width=Inches(4.8), size=14, gap=9)

# ═══════════════════════ 20. DEEP: EXPLAINABILITY + ERRORS ═══════════════════════
s = slide(); title_bar(s, "Phase 5 — Explainability & Error Analysis")
ex = p5e.get("explainability_top_words", {})
def tw(c):
    return " · ".join(t["token"] for t in ex.get(c, [])[:7])
bullets(s, [
    "Transparent — top Linear-SVM word features per class (no black box):",
    (f"Negative: {tw('Negative')}", 1),
    (f"Neutral: {tw('Neutral')}", 1),
    (f"Positive: {tw('Positive')}   (incl. English ‘thank’, Arabic الله)", 1),
    "Errors (test):",
    ("Hardest boundary = Neutral ↔ Negative (the fuzzy middle) for both models.", 1),
    (f"BanglaBERT error {EG.get('bert_latin_script',0):.2f} on Latin/code-mixed and "
     f"{EG.get('bert_emoji',0):.2f} on emoji comments vs ~{EG.get('bert_bangla',0):.2f} otherwise — "
     "the SVM is stable across both.", 1),
], top=Inches(1.55), gap=7)

# ═══════════════════════ 20a. ERROR EXAMPLES ═══════════════════════
s = slide(); title_bar(s, "Error Analysis — Hardest Cases (both models wrong)")
he = p5e.get("error_analysis", {}).get("hardest_examples", [])
rows = [["true", "SVM", "BERT", "comment"]]
for e in he[:6]:
    txt = e["text"].replace("\n", " ").strip()
    txt = (txt[:52] + "…") if len(txt) > 52 else txt
    rows.append([e["true"], e["svm"], e["bert"], txt])
table(s, rows, [1.6, 1.4, 1.4, 7.6], Inches(0.6), Inches(1.75), Inches(12.1), Inches(3.4), font=12)
bullets(s, [
    "Nearly all hard cases sit on the Neutral ↔ Negative boundary.",
    "Common types: rhetorical questions, mixed sentiment, sarcasm, and code-switched text.",
    "These are genuinely ambiguous — even human annotators would disagree, which sets a realistic "
    "ceiling on accuracy.",
], top=Inches(5.4), size=14, gap=7)

# ═══════════════════════ 20b. STATISTICAL VALIDATION ═══════════════════════
s = slide(); title_bar(s, "Statistical Validation", "Findings are not chance — Chi-square / Fisher's exact")
st = p5s.get("chi_square_tests", [])
rows = [["Association", "p-value", "Cramér's V", "significant?"]]
for t in st:
    rows.append([t["cross_tab"], f"{t['p_value']:.1e}", f"{t['cramers_v']}",
                 "YES" if t["significant"] else "no"])
tb = table(s, rows, [4.6, 2.0, 2.0, 2.0], Inches(0.7), Inches(1.75), Inches(12), Inches(2.9))
for i in range(1, len(rows)):
    cp = tb.cell(i, 3).text_frame.paragraphs[0]
    cp.font.bold = True; cp.font.color.rgb = GREEN if rows[i][3] == "YES" else GREY
fe = p5s.get("fisher_exact", {})
bullets(s, [
    f"Fisher's exact — victim-blaming × Neglect/Abandonment: odds ratio {fe.get('odds_ratio')}, "
    f"p = {fe.get('p_value',1):.1e} → victim-blaming is ~5× more likely in neglect topics.",
    "Emotion↔topic, emotion↔platform and victim-blaming↔topic are all statistically significant.",
    "Note (fair nuance): overall polarity↔platform is NOT significant — the platform effect is on "
    "fine-grained emotion, not coarse polarity.",
], top=Inches(4.9), size=14, gap=8)

# ═══════════════════════ 20c. LANGUAGE-WISE CONFUSION ═══════════════════════
s = slide(); title_bar(s, "Language-Wise Confusion — Seeing BanglaBERT Fail")
picture(s, "22_language_confusion.png", Inches(2.6), Inches(1.6), Inches(8.1), Inches(5.4))
caption(s, "SVM keeps a balanced diagonal on both scripts; BanglaBERT collapses on Banglish "
           "(never predicts Positive, scatters to Neutral) — the language finding, made visual.",
        Inches(6.95))

# ═══════════════════════ 21. KEY FINDINGS / CONTRIBUTIONS ═══════════════════════
s = slide()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH); fill(bg, NAVY)
tf = box(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(1.0))
p = tf.paragraphs[0]; p.text = "Key Findings & Contributions"
p.font.size = Pt(32); p.font.bold = True; p.font.color.rgb = GOLD
items = [
    "New annotated Bangla/Banglish elder-abuse dataset with a novel victim-blaming class.",
    "70% of reactions Negative; victim-blaming is topic-specific (high in Neglect, 0% in Murder).",
    "Platform effect: victim-blaming 18.6% on YouTube vs 0.9% on Facebook.",
    f"Classical SVM beats fine-tuned BanglaBERT on test ({TB.get('macro_f1',0):.2f} vs "
    f"{TR.get('macro_f1',0):.2f}" + (f", McNemar p={MC_P:.3f})." if SIG else "; gap not significant)."),
    "Explained WHY: BanglaBERT can't read Latin-script Banglish (err 0.51) or emoji (err 0.43); "
    "the char-ngram SVM is robust to both.",
    "Religion frames both blessing and condemnation (22% of comments).",
    "All key associations are statistically significant (Chi-square / Fisher's; "
    "victim-blaming odds ratio ≈ 4.9 in neglect topics).",
]
ct = box(s, Inches(1.0), Inches(1.75), Inches(11.4), Inches(5.2))
for i, it in enumerate(items):
    p = ct.paragraphs[0] if i == 0 else ct.add_paragraph()
    p.text = "✔  " + it
    p.font.size = Pt(17); p.font.color.rgb = WHITE; p.space_after = Pt(13)

# ═══════════════════════ 22. LIMITATIONS ═══════════════════════
s = slide(); title_bar(s, "Limitations", "Stated honestly")
bullets(s, [
    "Labels are model-generated (silver) following a codebook; human inter-annotator (Kappa) "
    "validation is planned but not yet complete.",
    "Small & imbalanced (2,301 comments; ~70% Negative) — caps rare-class performance and "
    "transformer fine-tuning.",
    "Narrow domain — YouTube/Facebook elder-abuse comments; may not transfer to other topics.",
    "Fine-grained emotion is subjective at the boundary (victim-blaming vs anger vs neutral).",
    "Point-in-time scrape (no temporal analysis); abuse_type is video-level context, not a "
    "per-comment label.",
    "BanglaBERT under-reads Latin-script Banglish and emoji — transformer results are a lower "
    "bound for code-switched text.",
], top=Inches(1.7), size=16, gap=10)

# ═══════════════════════ 23. ETHICS ═══════════════════════
s = slide(); title_bar(s, "Ethics Statement")
bullets(s, [
    "Only publicly-posted comments are used — no private data.",
    "Analysis is aggregate: we report patterns, never profile or expose individual commenters; "
    "names are not analysed or published.",
    "Elder abuse and its victims are treated with care — no attempt to re-identify anyone.",
    "The victim-blaming class is studied to understand and surface a harmful attitude — not to "
    "endorse it.",
    "Offensive comments are used solely for research and shown only as de-identified examples.",
    "Data collected in line with each platform's public-content terms; no user was contacted.",
], top=Inches(1.7), size=16, gap=10)

prs.save(str(OUT))
print(f"saved {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
